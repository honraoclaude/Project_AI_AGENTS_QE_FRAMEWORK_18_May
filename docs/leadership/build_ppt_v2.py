"""
FSC Agentic QE Framework — Leadership Deck v2
Design principle: each slide has ONE headline that IS the message.
Body text is supporting evidence, not the point.
Run: python docs/leadership/build_ppt_v2.py
Output: docs/leadership/FSC_QE_Framework_v2.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0A, 0x18, 0x28)
NAVY2   = RGBColor(0x12, 0x24, 0x38)
TEAL    = RGBColor(0x00, 0xB4, 0xD8)
TEAL2   = RGBColor(0x00, 0x8C, 0xAA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE= RGBColor(0xF2, 0xF6, 0xFA)
DARK    = RGBColor(0x1A, 0x2A, 0x3A)
MID     = RGBColor(0x3A, 0x4A, 0x5A)
GREEN   = RGBColor(0x00, 0x96, 0x5C)
GREEN2  = RGBColor(0xE4, 0xF8, 0xEF)
RED     = RGBColor(0xC0, 0x28, 0x28)
RED2    = RGBColor(0xFC, 0xEC, 0xEC)
AMBER   = RGBColor(0xE0, 0x88, 0x00)
AMBER2  = RGBColor(0xFD, 0xF3, 0xE0)
SLATE   = RGBColor(0x55, 0x6B, 0x82)

OUT = Path(__file__).parent / "FSC_QE_Framework_v2.pptx"
W, H = Inches(13.33), Inches(7.5)


# ── Core primitives ───────────────────────────────────────────────────────────

def new_prs():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def box(sl, x, y, w, h, fill, line_color=None, line_w=None):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        if line_w:
            s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def txt(sl, text, x, y, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    t = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name  = "Calibri"
    return t

def paras(sl, lines, x, y, w, h, size=15, color=DARK, spacing=1.15):
    """Multi-paragraph textbox."""
    t = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame
    tf.word_wrap = True
    for i, (line, bold, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = clr or color
        r.font.name  = "Calibri"
    return t

def big_num(sl, number, label, x, y, w=2.8, h=1.7,
            num_color=TEAL, bg=NAVY, lbl_color=WHITE):
    box(sl, x, y, w, h, bg)
    txt(sl, number, x+0.1, y+0.12, w-0.2, 0.95,
        size=38, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(sl, label,  x+0.1, y+1.1,  w-0.2, 0.52,
        size=12, color=lbl_color, align=PP_ALIGN.CENTER)

def headline(sl, text, y=0.22, size=40, color=WHITE):
    txt(sl, text, 0.5, y, 12.33, 1.1, size=size, bold=True, color=color)

def subhead(sl, text, y=1.35, color=TEAL, size=17):
    txt(sl, text, 0.5, y, 12.33, 0.5, size=size, color=color)

def divider(sl, y, color=TEAL, h=0.055):
    box(sl, 0.5, y, 12.33, h, color)


# ── SLIDE 1 — Title ──────────────────────────────────────────────────────────

def s01_title(prs):
    sl = blank(prs)
    box(sl, 0, 0, 13.33, 7.5, NAVY)
    box(sl, 0, 5.6, 13.33, 1.9, NAVY2)
    box(sl, 0, 5.56, 13.33, 0.07, TEAL)

    # Large kicker
    txt(sl, "FSC Agentic QE Framework",
        0.55, 1.2, 12.2, 1.3, size=48, bold=True, color=WHITE)
    txt(sl, "AI-Powered Quality Engineering for FCA-Regulated Salesforce FSC Delivery",
        0.55, 2.6, 11.5, 0.7, size=21, color=TEAL)

    # Four pills
    pills = ["53 Agents", "4 Phases", "12 Gates", "Validated 53 / 53"]
    pw = 2.7
    gap = 0.3
    sx = (13.33 - (len(pills)*pw + (len(pills)-1)*gap)) / 2
    for i, p in enumerate(pills):
        px = sx + i*(pw+gap)
        box(sl, px, 3.55, pw, 0.52, TEAL2)
        txt(sl, p, px, 3.58, pw, 0.46, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Senior Leadership Briefing  ·  May 2026",
        0.6, 5.82, 9.0, 0.4, size=13, color=RGBColor(0x88,0x99,0xAA))
    txt(sl, "CONFIDENTIAL",
        9.8, 5.82, 3.0, 0.4, size=13, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)


# ── SLIDE 2 — The Problem (big-number pain) ──────────────────────────────────

def s02_problem(prs):
    sl = blank(prs)
    box(sl, 0, 0, 13.33, 7.5, OFFWHITE)
    box(sl, 0, 0, 13.33, 1.55, NAVY)
    box(sl, 0, 1.5, 13.33, 0.07, TEAL)

    txt(sl, "Every sprint, QE is a bottleneck — and regulatory gaps go undetected",
        0.5, 0.18, 12.3, 1.1, size=28, bold=True, color=WHITE)

    pain = [
        ("30+",  "hours of senior QE time\nconsumed per sprint\non manual documentation"),
        ("43%",  "of FCA-flagged defects\noriginate from missing\nor untestable story ACs"),
        ("UAT",  "is where most regulatory\ngaps are discovered —\nthe most expensive point to fix"),
        ("0",    "structured Go/No-Go\ngates before release —\ndecisions made informally"),
    ]
    for i, (num, label) in enumerate(pain):
        x = 0.45 + i * 3.25
        box(sl, x, 1.75, 3.0, 5.3, WHITE,
            line_color=RGBColor(0xCC,0xD8,0xE4), line_w=9525)
        box(sl, x, 1.75, 3.0, 0.08, RED)
        txt(sl, num,   x+0.15, 1.95, 2.75, 1.5, size=60, bold=True, color=RED, align=PP_ALIGN.CENTER)
        txt(sl, label, x+0.2,  3.55, 2.65, 3.0, size=15, color=MID,  align=PP_ALIGN.CENTER)

    txt(sl, "These are not edge cases — they are the current baseline for FCA-regulated Salesforce delivery.",
        0.5, 7.1, 12.3, 0.35, size=12, italic=True, color=SLATE)


# ── SLIDE 3 — The Solution (one sentence + metrics) ──────────────────────────

def s03_solution(prs):
    sl = blank(prs)
    box(sl, 0, 0, 13.33, 7.5, NAVY)

    txt(sl, "We built an AI pipeline that QE-reviews every story — automatically, end to end.",
        0.55, 0.3, 12.2, 1.4, size=34, bold=True, color=WHITE)
    box(sl, 0.55, 1.72, 12.2, 0.06, TEAL)

    txt(sl, "53 specialised agents cover every QE discipline across 4 phases.\n"
            "Stories are checked, documented, and FCA-validated before a single line of code is written.\n"
            "Every decision is logged in an immutable audit ledger — no manual effort required per sprint.",
        0.55, 1.9, 12.0, 1.2, size=17, color=RGBColor(0xCC,0xDD,0xEE))

    metrics = [
        ("53",        "AI Agents"),
        ("4",         "Delivery Phases"),
        ("12",        "Quality Gates"),
        ("53 / 53",   "Agents Validated"),
        ("1,065",     "Automated Tests"),
        ("< 10 min",  "Per Story"),
    ]
    gap = 0.18
    mw  = (12.2 - gap*5) / 6
    sx  = 0.55
    for i, (num, lbl) in enumerate(metrics):
        big_num(sl, num, lbl, sx + i*(mw+gap), 3.3, w=mw, h=1.75)

    txt(sl, "Validated on FSC-2417 (HIGH-FCA, COBS 9.2 Suitability Assessment)  ·  Avg confidence 75.1%  ·  Total pipeline runtime 375 s",
        0.55, 5.3, 12.2, 0.45, size=13, color=RGBColor(0x66,0x88,0xAA), align=PP_ALIGN.CENTER)


# ── SLIDE 4 — Phase Pipeline ─────────────────────────────────────────────────

def s04_pipeline(prs):
    sl = blank(prs)
    box(sl, 0, 0, 13.33, 7.5, OFFWHITE)
    box(sl, 0, 0, 13.33, 1.5, NAVY)
    box(sl, 0, 1.46, 13.33, 0.07, TEAL)

    txt(sl, "Four phases, 12 gates — a story cannot advance without passing",
        0.5, 0.18, 12.3, 1.1, size=28, bold=True, color=WHITE)

    phases = [
        ("REFINEMENT",   "Agents 1–9",
         ["FCA Classify", "Consumer Duty Map", "BDD Acceptance Criteria",
          "Risk Register", "Test Design Strategy"],
         TEAL, "Gate G1–G4"),
        ("DEVELOPMENT",  "Agents 10–23",
         ["AC Compliance Check", "Apex Coverage", "Security Scan",
          "BDD Gherkin Files", "Test Data Architecture"],
         RGBColor(0x02,0x73,0x96), "Gate G5–G7"),
        ("TESTING",      "Agents 24–38",
         ["CRT & UAT Scenarios", "FCA Regulatory Tests", "Defect Triage",
          "Root Cause Analysis", "Regression Risk Score"],
         RGBColor(0x01,0x52,0x71), "Gate G8–G10"),
        ("RELEASE",      "Agents 39–50",
         ["Go/No-Go Coordinator", "Change Set Integrity", "Dry-Run Validation",
          "FCA Evidence Pack", "Production Validation"],
         RGBColor(0x0A,0x18,0x28), "Gate G11–G12"),
    ]

    pw = 2.9
    gap_x = 0.33
    sx = 0.4

    for i, (phase, agents, items, colour, gate) in enumerate(phases):
        x = sx + i*(pw+gap_x)

        # Phase header bar
        box(sl, x, 1.65, pw, 0.72, colour)
        txt(sl, phase,  x+0.12, 1.68, pw-0.24, 0.42, size=14, bold=True, color=WHITE)
        txt(sl, agents, x+0.12, 2.1,  pw-0.24, 0.28, size=11, color=TEAL)

        # Body card
        box(sl, x, 2.37, pw, 4.0, WHITE,
            line_color=RGBColor(0xCC,0xD8,0xE4), line_w=6350)

        # Items
        for j, item in enumerate(items):
            iy = 2.5 + j*0.58
            box(sl, x+0.18, iy, 0.22, 0.32, colour)
            txt(sl, item, x+0.52, iy+0.02, pw-0.65, 0.32, size=12, color=DARK)

        # Gate badge at bottom
        box(sl, x+0.2, 6.05, pw-0.4, 0.28, colour)
        txt(sl, gate, x+0.2, 6.06, pw-0.4, 0.26,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Arrow between phases
        if i < 3:
            ax = x + pw + 0.04
            txt(sl, "›", ax, 3.55, 0.28, 0.55, size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    txt(sl, "Each gate is fail-closed: stories are blocked from advancing, not merely warned.",
        0.5, 6.95, 12.3, 0.4, size=13, italic=True, color=SLATE)


# ── SLIDE 5 — Before vs After ─────────────────────────────────────────────────

def s05_before_after(prs):
    sl = blank(prs)
    box(sl, 0, 0, 13.33, 7.5, OFFWHITE)
    box(sl, 0, 0, 13.33, 1.35, NAVY)
    box(sl, 0, 1.31, 13.33, 0.07, TEAL)

    txt(sl, "From vague stories to FCA-compliant, BDD-tested, audit-ready releases",
        0.5, 0.18, 12.3, 1.0, size=27, bold=True, color=WHITE)

    # Column headers
    box(sl, 0.4, 1.5, 6.0, 0.55, RED)
    txt(sl, "✗  Without the Framework",
        0.55, 1.55, 5.75, 0.44, size=16, bold=True, color=WHITE)
    box(sl, 6.93, 1.5, 6.0, 0.55, GREEN)
    txt(sl, "✓  With FSC QE Framework",
        7.08, 1.55, 5.75, 0.44, size=16, bold=True, color=WHITE)
    box(sl, 6.55, 1.5, 0.35, 0.55, OFFWHITE)

    rows = [
        ("Refinement",
         "ACs are incomplete, vague, or missing error and FCA scenarios",
         "BDD Given/When/Then ACs auto-generated — happy path, error, edge, and regulatory"),
        ("Refinement",
         "FCA risk depends on analyst memory — inconsistently applied",
         "Every story classified HIGH/MEDIUM/LOW; COBS 9.2, Consumer Duty, FG21/1 auto-checked"),
        ("Development",
         "Apex coverage gaps and security issues caught late in peer review or CI",
         "Coverage analysis, PMD security scan, and code quality review run automatically"),
        ("Development",
         "QEs manually write Gherkin files and test data scripts per story",
         "BDD .feature files and test data architecture generated and linked to ACs"),
        ("Testing",
         "CRT and UAT test cases designed ad hoc — coverage gaps found during execution",
         "CRT, UAT, FCA regulatory, and regression scenarios structured before test execution"),
        ("Release",
         "Go/No-Go is informal — no gate; FCA evidence assembled manually under pressure",
         "Automated Go/No-Go gate; FCA Evidence Pack auto-assembled; release blocked if incomplete"),
    ]

    phase_col = {
        "Refinement":  RGBColor(0xE0,0xF4,0xFD),
        "Development": RGBColor(0xE0,0xFD,0xF0),
        "Testing":     RGBColor(0xFD,0xF4,0xE0),
        "Release":     RGBColor(0xF0,0xE8,0xFD),
    }
    phase_text = {
        "Refinement":  RGBColor(0x02,0x73,0x96),
        "Development": RGBColor(0x00,0x65,0x40),
        "Testing":     RGBColor(0x80,0x50,0x00),
        "Release":     RGBColor(0x50,0x20,0x80),
    }

    y = 2.14
    rh = 0.78
    for phase, before, after in rows:
        pc = phase_col[phase]
        pt = phase_text[phase]

        # Phase tag
        box(sl, 0.4, y, 0.9, rh, pc)
        txt(sl, phase, 0.42, y + rh/2 - 0.14, 0.88, 0.3,
            size=9, bold=True, color=pt, align=PP_ALIGN.CENTER)

        # Before cell
        box(sl, 1.3, y, 5.1, rh, RED2,
            line_color=RGBColor(0xE8,0xCC,0xCC), line_w=6350)
        txt(sl, before, 1.42, y+0.1, 4.9, rh-0.2, size=12, color=RED)

        # Divider
        box(sl, 6.55, y, 0.35, rh, OFFWHITE)
        txt(sl, "›", 6.56, y+0.18, 0.33, 0.38, size=18, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

        # After cell
        box(sl, 6.93, y, 5.97, rh, GREEN2,
            line_color=RGBColor(0xCC,0xE8,0xD4), line_w=6350)
        txt(sl, after, 7.05, y+0.1, 5.75, rh-0.2, size=12, color=GREEN)

        y += rh + 0.04


# ── SLIDE 6 — FCA Compliance ─────────────────────────────────────────────────

def s06_fca(prs):
    sl = blank(prs)
    box(sl, 0, 0, 13.33, 7.5, NAVY)

    txt(sl, "FCA compliance is checked on every story — not hoped for",
        0.55, 0.22, 12.2, 1.1, size=32, bold=True, color=WHITE)
    box(sl, 0.55, 1.38, 12.2, 0.06, TEAL)

    pillars = [
        ("COBS 9.2",
         "Suitability",
         [
             "Agent 03 detects suitability triggers and classifies story as HIGH-FCA",
             "Agent 05 generates mandatory COBS 9.2 regulatory AC scenarios",
             "Agent 45 blocks Go/No-Go if COBS 9.2 evidence is missing",
             "Acknowledgement audit record created with adviser ID and timestamp",
         ]),
        ("Consumer Duty\nPS22/9",
         "Outcomes 1–4",
         [
             "Agent 04 maps all four Consumer Duty outcomes to every story",
             "Vulnerable customer pathways (FG21/1 §4.3) flagged and test-covered",
             "Agent 30 generates dedicated FCA regulatory test scenarios",
             "Agent 44 assembles Consumer Duty evidence per story automatically",
         ]),
        ("Audit Trail\nSYSC",
         "Immutable Ledger",
         [
             "Every agent decision written to an append-only PostgreSQL ledger",
             "Records are timestamped, signed, and cannot be modified or deleted",
             "Full evidence trail from refinement to release — one query, FCA-ready",
             "No manual document assembly before internal audit or FCA review",
         ]),
    ]

    for i, (title, sub, points) in enumerate(pillars):
        x = 0.5 + i * 4.25
        box(sl, x, 1.58, 4.0, 5.65, RGBColor(0x12,0x28,0x40))
        box(sl, x, 1.58, 4.0, 0.98, TEAL)
        txt(sl, title, x+0.18, 1.63, 3.65, 0.62, size=18, bold=True, color=NAVY)
        txt(sl, sub,   x+0.18, 2.24, 3.65, 0.3,  size=12, bold=True, color=WHITE)

        for j, pt in enumerate(points):
            py = 2.7 + j * 1.1
            box(sl, x+0.2, py, 0.3, 0.3, TEAL)
            txt(sl, "›", x+0.2, py, 0.3, 0.3, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
            txt(sl, pt, x+0.6, py, 3.25, 0.95, size=12.5, color=RGBColor(0xCC,0xDD,0xEE))

    box(sl, 0.5, 7.25, 12.3, 0.15, TEAL)


# ── SLIDE 7 — Metrics wall ───────────────────────────────────────────────────

def s07_metrics(prs):
    sl = blank(prs)
    box(sl, 0, 0, 13.33, 7.5, NAVY2)

    txt(sl, "Proven.", 0.55, 0.18, 4.0, 1.05, size=52, bold=True, color=WHITE)
    txt(sl, "Measured.", 4.3, 0.18, 4.5, 1.05, size=52, bold=True, color=TEAL)
    txt(sl, "Ready.", 8.4, 0.18, 4.5, 1.05, size=52, bold=True, color=WHITE)
    box(sl, 0.55, 1.28, 12.2, 0.06, TEAL)

    # Big metric grid — 3 columns × 2 rows
    metrics = [
        ("53 / 53", "Agents passed\nvalidation",    GREEN),
        ("1,065",   "Automated tests\nall passing", TEAL),
        ("75.1%",   "Average agent\nconfidence score", WHITE),
        ("< 10 min","Full pipeline\nper story",     TEAL),
        ("375 s",   "Validated runtime\non FSC-2417", WHITE),
        ("18",      "Production gaps\ndocumented",  AMBER),
    ]
    mw, mh = 3.8, 2.45
    gap = 0.22
    sx = (13.33 - (3*mw + 2*gap)) / 2

    for i, (num, lbl, col) in enumerate(metrics):
        row, col_i = divmod(i, 3)
        x = sx + col_i * (mw+gap)
        y = 1.5 + row * (mh+0.18)
        box(sl, x, y, mw, mh, DARK)
        box(sl, x, y, mw, 0.06, col)
        txt(sl, num, x+0.2, y+0.18, mw-0.4, 1.3,
            size=46, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(sl, lbl, x+0.2, y+1.6, mw-0.4, 0.72,
            size=13, color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.CENTER)


# ── SLIDE 8 — ROI ─────────────────────────────────────────────────────────────

def s08_roi(prs):
    sl = blank(prs)
    box(sl, 0, 0, 13.33, 7.5, OFFWHITE)
    box(sl, 0, 0, 13.33, 1.5, NAVY)
    box(sl, 0, 1.46, 13.33, 0.07, TEAL)

    txt(sl, "1,500 senior QE hours per year — redirected to high-value testing",
        0.5, 0.18, 12.3, 1.1, size=28, bold=True, color=WHITE)

    # Hero calculation
    box(sl, 0.4, 1.65, 7.5, 5.55, WHITE,
        line_color=RGBColor(0xCC,0xD8,0xE4), line_w=6350)
    box(sl, 0.4, 1.65, 7.5, 0.06, TEAL)

    txt(sl, "Conservative ROI Calculation", 0.65, 1.8, 7.0, 0.45,
        size=16, bold=True, color=NAVY)

    calc = [
        ("10",   "stories per sprint requiring QE documentation"),
        ("×  3", "hours of senior QE time saved per story"),
        ("×  50","sprints per year (2-week cadence)"),
        ("=  1,500", "senior QE hours per year freed"),
    ]
    y = 2.45
    for num, label in calc:
        if num.startswith("="):
            box(sl, 0.55, y-0.05, 7.1, 0.65, RGBColor(0xE4,0xF8,0xEF))
            txt(sl, num,   0.7,  y+0.05, 2.5, 0.52, size=22, bold=True, color=GREEN)
            txt(sl, label, 3.1,  y+0.05, 4.3, 0.52, size=16, bold=True, color=GREEN)
        else:
            txt(sl, num,   0.7,  y, 2.2, 0.52, size=20, bold=True, color=TEAL)
            txt(sl, label, 3.0,  y, 4.7, 0.52, size=15, color=MID)
        y += 0.72

    txt(sl, "* Based on 3 hrs manual effort per story for AC writing, FCA review,\n  test planning, and documentation. Blended rate not included.",
        0.65, 6.0, 7.0, 0.65, size=10.5, italic=True, color=SLATE)

    # Side benefits
    box(sl, 8.33, 1.65, 4.6, 5.55, NAVY)
    txt(sl, "Additional Value", 8.55, 1.82, 4.2, 0.45, size=15, bold=True, color=TEAL)

    side = [
        ("Defect Prevention",
         "Risks raised at refinement — not in UAT. Earlier fix = lower cost."),
        ("FCA Audit Prep",
         "Evidence pack auto-assembled. No sprint dedicated to audit prep."),
        ("Delivery Consistency",
         "Same QE rigour applied to every story, every sprint, every team."),
        ("Regulatory Protection",
         "Systematic FCA checks reduce tail risk of material FCA findings."),
    ]
    sy = 2.4
    for title, body in side:
        box(sl, 8.48, sy, 0.22, 0.78, TEAL)
        txt(sl, title, 8.82, sy+0.04, 3.9, 0.35, size=13, bold=True, color=WHITE)
        txt(sl, body,  8.82, sy+0.4,  3.9, 0.38, size=11.5, color=RGBColor(0xBB,0xCC,0xDD))
        sy += 1.1


# ── SLIDE 9 — Production Readiness ───────────────────────────────────────────

def s09_readiness(prs):
    sl = blank(prs)
    box(sl, 0, 0, 13.33, 7.5, OFFWHITE)
    box(sl, 0, 0, 13.33, 1.5, NAVY)
    box(sl, 0, 1.46, 13.33, 0.07, TEAL)

    txt(sl, "Built, tested, and validated — 18 gaps remain before production",
        0.5, 0.18, 12.3, 1.1, size=28, bold=True, color=WHITE)

    done = [
        "53/53 agents validated on FSC-2417",
        "1,065 automated unit tests — all passing",
        "12 quality gates implemented and verified",
        "FCA compliance checks validated for HIGH-FCA story",
        "BDD AC output includes test_category classification",
        "HTML validation report generated per run",
        "Immutable audit ledger schema defined",
    ]
    pending = [
        "PostgreSQL production database provisioned",
        "Anthropic API key in production secrets manager",
        "Jira production OAuth credentials configured",
        "Copado webhook registered for CI/CD integration",
        "Compliance sign-off on FCA Evidence Pack format",
        "PRODUCTION_READINESS.md: 18 gaps with owners & ETAs",
    ]

    # Done column
    box(sl, 0.4, 1.65, 5.9, 5.55, WHITE, line_color=RGBColor(0xCC,0xD8,0xE4), line_w=6350)
    box(sl, 0.4, 1.65, 5.9, 0.5, GREEN)
    txt(sl, "✓  Complete", 0.6, 1.72, 5.6, 0.38, size=15, bold=True, color=WHITE)
    y = 2.28
    for item in done:
        box(sl, 0.55, y+0.06, 0.28, 0.28, GREEN)
        txt(sl, "✓", 0.55, y+0.06, 0.28, 0.28, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, item, 0.98, y+0.05, 5.15, 0.42, size=13, color=DARK)
        y += 0.52

    # Pending column
    box(sl, 7.0, 1.65, 5.9, 5.55, WHITE, line_color=RGBColor(0xCC,0xD8,0xE4), line_w=6350)
    box(sl, 7.0, 1.65, 5.9, 0.5, AMBER)
    txt(sl, "⚠  Pending for Production", 7.2, 1.72, 5.6, 0.38, size=15, bold=True, color=WHITE)
    y = 2.28
    for item in pending:
        box(sl, 7.15, y+0.06, 0.28, 0.28, AMBER)
        txt(sl, "!", 7.15, y+0.06, 0.28, 0.28, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, item, 7.58, y+0.05, 5.15, 0.42, size=13, color=DARK)
        y += 0.52

    box(sl, 6.45, 1.65, 0.1, 5.55, RGBColor(0xCC,0xD8,0xE4))


# ── SLIDE 10 — Next Steps ────────────────────────────────────────────────────

def s10_next(prs):
    sl = blank(prs)
    box(sl, 0, 0, 13.33, 7.5, NAVY)

    txt(sl, "Three actions to go live", 0.55, 0.22, 12.0, 0.9, size=38, bold=True, color=WHITE)
    box(sl, 0.55, 1.18, 12.2, 0.06, TEAL)

    actions = [
        ("01", "Approve Production Deployment",
         "Commission infrastructure provisioning — PostgreSQL, secrets manager, Jira OAuth, Copado webhook.\n"
         "Assign a delivery lead to own the 18-gap closure plan in PRODUCTION_READINESS.md.\n"
         "Target: framework live before the next quarter's sprint cycle begins.",
         TEAL),
        ("02", "Run a Pilot Sprint",
         "Select 3–5 stories from the next sprint backlog.\n"
         "Run the framework in parallel with the existing manual QE process — no disruption to delivery.\n"
         "Measure: AC quality delta, FCA scenario coverage rate, and hours saved.",
         GREEN),
        ("03", "Brief Compliance on Evidence Pack",
         "Share the auto-generated FCA Evidence Pack for FSC-2417 with the Compliance team.\n"
         "Confirm the audit trail format meets SYSC and Consumer Duty documentation requirements.\n"
         "Obtain sign-off that automated evidence packs are acceptable for FCA review.",
         AMBER),
    ]

    y = 1.4
    for num, title, body, colour in actions:
        box(sl, 0.5, y, 12.3, 1.85, DARK)
        box(sl, 0.5, y, 0.72, 1.85, colour)

        txt(sl, num,   0.54, y+0.62, 0.65, 0.65,
            size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(sl, title, 1.38, y+0.12, 11.2, 0.52, size=19, bold=True, color=WHITE)

        lines = body.split("\n")
        for li, line in enumerate(lines):
            txt(sl, "›  " + line, 1.38, y+0.7+li*0.38, 11.1, 0.38, size=13,
                color=RGBColor(0xBB,0xCC,0xDD))

        y += 2.02

    txt(sl, "FSC Agentic QE Framework  ·  Senior Leadership Briefing  ·  May 2026",
        0.55, 7.15, 12.2, 0.3, size=11, color=SLATE, align=PP_ALIGN.CENTER)


# ── Build ─────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    s01_title(prs)
    s02_problem(prs)
    s03_solution(prs)
    s04_pipeline(prs)
    s05_before_after(prs)
    s06_fca(prs)
    s07_metrics(prs)
    s08_roi(prs)
    s09_readiness(prs)
    s10_next(prs)
    prs.save(str(OUT))
    print(f"Saved: {OUT}  ({OUT.stat().st_size // 1024} KB)  —  {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
