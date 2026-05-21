"""
Build FSC Agentic QE Framework — Leadership PowerPoint
Run: python docs/leadership/build_ppt.py
Output: docs/leadership/FSC_QE_Framework_Leadership.pptx
No API calls — pure python-pptx.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # title slides, headers
TEAL   = RGBColor(0x00, 0xB4, 0xD8)   # accent, highlights
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # content slide background
DARK   = RGBColor(0x1A, 0x2A, 0x3A)   # body text on light bg
GREEN  = RGBColor(0x00, 0x96, 0x5C)   # success / validated
AMBER  = RGBColor(0xE8, 0x8C, 0x00)   # caution

OUT_PATH = Path(__file__).parent / "FSC_QE_Framework_Leadership.pptx"
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # completely blank


def rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def textbox(slide, text, x, y, w, h,
            size=18, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb


def multiline(slide, lines, x, y, w, h,
              size=16, color=DARK, bold_first=False):
    """Add a textbox with multiple paragraphs."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        run.font.bold  = (bold_first and i == 0)
        run.font.name  = "Calibri"
    return txb


def header_bar(slide, title, subtitle=None):
    """Navy top bar with white title."""
    rect(slide, 0, 0, 13.33, 1.35, NAVY)
    textbox(slide, title,   0.4, 0.1, 12.5, 0.75,
            size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        textbox(slide, subtitle, 0.4, 0.82, 12.5, 0.4,
                size=16, color=TEAL, align=PP_ALIGN.LEFT)


def teal_pill(slide, text, x, y, w=2.6, h=0.45):
    """Small teal label pill."""
    rect(slide, x, y, w, h, TEAL)
    textbox(slide, text, x + 0.08, y + 0.03, w - 0.16, h - 0.06,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def metric_box(slide, number, label, x, y):
    """Big number + label metric card."""
    rect(slide, x, y, 2.8, 1.6, NAVY)
    textbox(slide, number, x + 0.1, y + 0.1, 2.6, 0.9,
            size=36, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    textbox(slide, label,  x + 0.1, y + 1.0, 2.6, 0.55,
            size=13, color=WHITE, align=PP_ALIGN.CENTER)


def section_label(slide, text, x=0.4, y=1.5):
    textbox(slide, text, x, y, 12.5, 0.4,
            size=13, bold=True, color=TEAL, align=PP_ALIGN.LEFT)


# ── SLIDE 1 — Title ──────────────────────────────────────────────────────────

def slide_title(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, NAVY)                   # full navy bg
    rect(sl, 0, 5.5, 13.33, 2.0, DARK)                  # darker footer strip
    rect(sl, 0, 5.45, 13.33, 0.08, TEAL)                # teal divider

    textbox(sl, "FSC Agentic QE Framework",
            0.7, 1.6, 12.0, 1.2,
            size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    textbox(sl, "AI-Powered Quality Engineering for FCA-Regulated Salesforce Delivery",
            0.7, 2.85, 11.5, 0.7,
            size=22, color=TEAL, align=PP_ALIGN.LEFT)
    textbox(sl, "53 Agents  ·  4 Phases  ·  12 Gates  ·  FCA-Aware  ·  Validated 53/53",
            0.7, 3.65, 11.5, 0.5,
            size=16, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.LEFT)

    textbox(sl, "Senior Leadership Briefing  |  May 2026",
            0.7, 5.75, 8.0, 0.4,
            size=14, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.LEFT)
    textbox(sl, "CONFIDENTIAL",
            9.5, 5.75, 3.5, 0.4,
            size=14, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)


# ── SLIDE 2 — The Problem ─────────────────────────────────────────────────────

def slide_problem(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT)
    header_bar(sl, "The Problem", "Manual QE in regulated delivery is slow, inconsistent, and risky")

    problems = [
        ("01", "FCA Obligations Are Missed",
         "COBS 9.2, Consumer Duty PS22/9, and FG21/1 checks rely on analyst memory.\n"
         "Regulatory scenarios are skipped under sprint pressure — creating audit exposure."),
        ("02", "Acceptance Criteria Quality Is Inconsistent",
         "Stories enter development with vague, incomplete, or untestable ACs.\n"
         "Defects caught in UAT or production that should have been specified at refinement."),
        ("03", "QE Is a Bottleneck, Not a Safeguard",
         "Manual test design, FCA review, and audit documentation consume 30+ hours\n"
         "per sprint of senior QE time — effort that does not scale with delivery pace."),
    ]

    for i, (num, title, body) in enumerate(problems):
        x = 0.4 + i * 4.3
        rect(sl, x, 1.6, 4.0, 5.3, WHITE)
        rect(sl, x, 1.6, 4.0, 0.55, TEAL)
        textbox(sl, num,   x + 0.15, 1.65, 0.5, 0.45, size=18, bold=True, color=WHITE)
        textbox(sl, title, x + 0.65, 1.65, 3.2, 0.45, size=14, bold=True, color=WHITE)
        multiline(sl, body.split("\n"), x + 0.2, 2.35, 3.65, 4.1, size=14, color=DARK)


# ── SLIDE 3 — The Solution ────────────────────────────────────────────────────

def slide_solution(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT)
    header_bar(sl, "The Solution", "An autonomous AI pipeline that QE-reviews every story — before a line of code is written")

    textbox(sl, "FSC Agentic QE Framework", 0.4, 1.55, 12.5, 0.55,
            size=22, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    points = [
        "53 specialised AI agents — each focused on one QE discipline",
        "4 sequential phases: Refinement → Development → Testing → Release",
        "12 quality gates — stories cannot advance if critical criteria fail",
        "FCA-aware — COBS 9.2, Consumer Duty PS22/9, FG21/1 checked automatically on every story",
        "Immutable audit ledger — every agent decision is logged and cannot be altered",
        "Generates structured BDD acceptance criteria, test design strategy, risk register, and FCA evidence pack",
        "Runs in under 10 minutes per story — no manual effort required per sprint",
    ]
    lines = ["  ›  " + p for p in points]
    multiline(sl, lines, 0.6, 2.2, 11.8, 4.2, size=17, color=DARK)

    rect(sl, 0.4, 6.4, 12.5, 0.75, NAVY)
    textbox(sl, "Result: every story exits refinement with complete, FCA-validated QE documentation — automatically.",
            0.7, 6.5, 12.0, 0.55, size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


# ── SLIDE 4 — How It Works ────────────────────────────────────────────────────

def slide_pipeline(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT)
    header_bar(sl, "How It Works", "Four phases, each gate-controlled — stories advance only when quality criteria are met")

    phases = [
        ("REFINEMENT", "Agents 1–9", ["Story Intent", "FCA Classify", "Consumer Duty",
                                       "AC Generator", "Test Design", "Risk Register"]),
        ("DEVELOPMENT", "Agents 10–23", ["AC Compliance", "Branch Tracer", "Apex Coverage",
                                          "Code Quality", "BDD Gherkin", "Test Data"]),
        ("TESTING", "Agents 24–38", ["CRT Scenarios", "UAT Cases", "FCA Scenarios",
                                      "Root Cause", "Regression Risk", "Flaky Tests"]),
        ("RELEASE", "Agents 39–50", ["Readiness", "Change Set", "Dry-Run", "FCA Evidence",
                                      "Go/No-Go", "Prod Validation"]),
    ]

    colours = [TEAL, RGBColor(0x02, 0x73, 0x96), RGBColor(0x01, 0x52, 0x71), NAVY]
    x_start = 0.35

    for i, (phase, agents, items) in enumerate(phases):
        x = x_start + i * 3.25
        rect(sl, x, 1.55, 3.0, 0.55, colours[i])
        textbox(sl, phase,  x + 0.1, 1.6, 2.8, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(sl, agents, x + 0.1, 2.15, 2.8, 0.35, size=12, color=colours[i], align=PP_ALIGN.CENTER)

        rect(sl, x, 2.55, 3.0, 4.1, WHITE)
        item_lines = ["  ✓  " + it for it in items]
        multiline(sl, item_lines, x + 0.15, 2.65, 2.75, 3.85, size=13, color=DARK)

        if i < 3:
            textbox(sl, "→", x + 3.06, 2.9, 0.25, 0.5, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    textbox(sl, "Each phase concludes with automated gate checks. Stories are blocked from advancing on failure — not just flagged.",
            0.4, 6.85, 12.5, 0.45, size=13, color=DARK, align=PP_ALIGN.LEFT)


# ── SLIDE 5 — Agent Capabilities ─────────────────────────────────────────────

def slide_capabilities(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT)
    header_bar(sl, "What the Framework Produces", "Outputs per story, automatically — replacing hours of manual QE work")

    rows = [
        ("Refinement",   "Complete BDD acceptance criteria (Given/When/Then) for every scenario type",    "Agent 05 AC Generator"),
        ("Refinement",   "FCA risk classification: HIGH / MEDIUM / LOW with regulatory triggers",          "Agent 03 FCA Classifier"),
        ("Refinement",   "Consumer Duty mapping — all PS22/9 obligations linked to story ACs",            "Agent 04 Consumer Duty"),
        ("Refinement",   "Risk register — critical/high risks flagged before development starts",          "Agent 09 Risk Anticipation"),
        ("Development",  "BDD Gherkin test files (.feature) ready for execution",                          "Agent 19 BDD Writer"),
        ("Development",  "Test data architecture — Salesforce setup scripts and data volumes",              "Agent 21 Test Data"),
        ("Testing",      "CRT scenarios, UAT test cases, FCA-specific regulatory tests",                   "Agents 26-30"),
        ("Release",      "FCA Evidence Pack — audit-ready compliance documentation",                       "Agent 44"),
        ("Release",      "Go/No-Go coordinator — final release gate with sign-off tracking",               "Agent 45"),
    ]

    col_widths = [1.9, 6.8, 2.8]
    col_x = [0.35, 2.35, 9.25]
    headers = ["Phase", "Output", "Agent"]

    y = 1.55
    for j, (h, cw, cx) in enumerate(zip(headers, col_widths, col_x)):
        rect(sl, cx, y, cw, 0.38, NAVY)
        textbox(sl, h, cx + 0.08, y + 0.04, cw - 0.16, 0.32,
                size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    phase_colours = {
        "Refinement":  RGBColor(0xE8, 0xF4, 0xFD),
        "Development": RGBColor(0xE8, 0xFD, 0xF4),
        "Testing":     RGBColor(0xFD, 0xF4, 0xE8),
        "Release":     RGBColor(0xF4, 0xE8, 0xFD),
    }

    y = 1.93
    for row_phase, row_output, row_agent in rows:
        row_h = 0.5
        bg = phase_colours.get(row_phase, WHITE)
        for cw, cx in zip(col_widths, col_x):
            rect(sl, cx, y, cw, row_h, bg)
        for text, cw, cx in zip([row_phase, row_output, row_agent], col_widths, col_x):
            textbox(sl, text, cx + 0.08, y + 0.05, cw - 0.16, row_h - 0.1,
                    size=11, color=DARK, align=PP_ALIGN.LEFT)
        y += row_h

    textbox(sl, "Each output is structured JSON stored in the audit ledger — queryable, reportable, and FCA-submittable.",
            0.35, 6.85, 12.5, 0.45, size=12, color=DARK)


# ── SLIDE 6 — FCA Compliance ──────────────────────────────────────────────────

def slide_fca(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, NAVY)
    rect(sl, 0, 1.4, 13.33, 6.1, LIGHT)
    header_bar(sl, "FCA Compliance — Built In, Not Bolted On", "Every story is checked against FCA obligations before development begins")

    regs = [
        ("COBS 9.2\nSuitability",
         "Agent 03 detects Suitability__c and RiskProfile__c involvement.\n"
         "Agent 05 generates mandatory regulatory AC scenarios testing the COBS 9.2 gate.\n"
         "Agent 45 blocks go/no-go if COBS 9.2 scenarios are missing."),
        ("Consumer Duty\nPS22/9",
         "Agent 04 maps all four Consumer Duty outcomes to story obligations.\n"
         "Vulnerable customer pathways (FG21/1 §4.3) are flagged and tested.\n"
         "Agent 44 assembles Consumer Duty evidence for each story automatically."),
        ("Audit Trail\n& SYSC",
         "All agent decisions written to an immutable, append-only PostgreSQL ledger.\n"
         "COBS 9.2 acknowledgements are captured with adviser ID and timestamp.\n"
         "Evidence pack is audit-ready — no manual assembly before FCA review."),
    ]

    for i, (title, body) in enumerate(regs):
        x = 0.4 + i * 4.3
        rect(sl, x, 1.6, 4.0, 5.4, WHITE)
        rect(sl, x, 1.6, 4.0, 0.85, TEAL)
        textbox(sl, title, x + 0.15, 1.65, 3.75, 0.75,
                size=16, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        multiline(sl, body.split("\n"), x + 0.2, 2.6, 3.65, 4.1, size=13, color=DARK)

    rect(sl, 0.4, 7.08, 12.5, 0.32, NAVY)
    textbox(sl,
            "Framework FCA classification is independently validated per story — HIGH, MEDIUM, or LOW — "
            "and gates are tighter for HIGH-FCA stories.",
            0.6, 7.1, 12.1, 0.28, size=11, color=WHITE)


# ── SLIDE 7 — Metrics ─────────────────────────────────────────────────────────

def slide_metrics(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, NAVY)

    textbox(sl, "Validated. Measured. Ready.",
            0.5, 0.3, 12.0, 0.8, size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    textbox(sl, "Framework validation results — story FSC-2417, Salesforce FSC Wealth Management",
            0.5, 1.1, 12.0, 0.45, size=16, color=TEAL, align=PP_ALIGN.LEFT)

    metrics = [
        ("53 / 53",   "Agents Passed"),
        ("1,065",     "Automated Tests"),
        ("75.1 %",    "Avg Confidence"),
        ("< 10 min",  "Per Story Runtime"),
    ]

    for i, (num, label) in enumerate(metrics):
        x = 0.5 + i * 3.25
        metric_box(sl, num, label, x, 1.85)

    details = [
        ("Confidence Tiers", "Tier A = 97 (deterministic) · Tier B = signal-scored (20–92) · Tier C = rule-based"),
        ("Escalation",       "Stories scoring below 60 are automatically flagged for senior QE review"),
        ("Coverage",         "Happy path · error paths · edge cases · regulatory scenarios — all verified per story"),
        ("Audit Ledger",     "Immutable append-only PostgreSQL log — every agent verdict stored with timestamp"),
        ("Gate Results",     "12 gates across 4 phases · stories blocked (not warned) on gate failure"),
    ]

    y = 3.75
    for label, detail in details:
        rect(sl, 0.5, y, 12.3, 0.52, DARK)
        textbox(sl, label,  0.65, y + 0.08, 2.2, 0.38, size=13, bold=True, color=TEAL)
        textbox(sl, detail, 2.9,  y + 0.08, 9.8, 0.38, size=13, color=WHITE)
        y += 0.57

    textbox(sl, "Validation run: 19 May 2026  ·  Story: FSC-2417 (HIGH-FCA, COBS 9.2 Suitability Assessment)",
            0.5, 7.1, 12.3, 0.35, size=11, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.LEFT)


# ── SLIDE 8 — Business Benefits ──────────────────────────────────────────────

def slide_benefits(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT)
    header_bar(sl, "Business Benefits", "Measurable impact across QE efficiency, regulatory risk, and delivery quality")

    benefits = [
        ("Time Saved",          TEAL,
         "Replaces 30+ hours of manual QE work per sprint.\n"
         "AC writing, FCA review, test design, and documentation\n"
         "are produced automatically in under 10 minutes."),
        ("Regulatory Risk",     GREEN,
         "FCA obligations checked on every story, every sprint.\n"
         "No dependency on analyst memory or checklist discipline.\n"
         "Audit evidence generated — not assembled under pressure."),
        ("Defect Prevention",   RGBColor(0x5B, 0x2D, 0x8E),
         "Risk register raised at refinement — before code is written.\n"
         "Critical risks are visible to the team 2–3 weeks earlier.\n"
         "BDD scenarios defined before development reduces rework."),
        ("Delivery Confidence", NAVY,
         "Go/No-Go gate with structured sign-off before every release.\n"
         "FCA Evidence Pack assembled automatically per story.\n"
         "Immutable audit ledger supports FCA and internal audit."),
    ]

    for i, (title, colour, body) in enumerate(benefits):
        x = 0.4 if i < 2 else 7.0
        y = 1.6 if i % 2 == 0 else 4.15
        w = 5.9
        rect(sl, x, y, w, 2.35, WHITE)
        rect(sl, x, y, 0.22, 2.35, colour)
        textbox(sl, title, x + 0.35, y + 0.12, w - 0.5, 0.42, size=16, bold=True, color=NAVY)
        multiline(sl, body.split("\n"), x + 0.35, y + 0.6, w - 0.5, 1.6, size=13, color=DARK)

    rect(sl, 0.4, 6.7, 12.5, 0.6, NAVY)
    textbox(sl,
            "Conservative ROI estimate: 10 stories/sprint × 3 hrs manual QE saved × 50 sprints/year "
            "= 1,500 senior QE hours per year redirected to high-value testing.",
            0.65, 6.78, 12.1, 0.45, size=13, color=WHITE)


# ── SLIDE 9 — Production Readiness ───────────────────────────────────────────

def slide_production(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, LIGHT)
    header_bar(sl, "Production Readiness", "Framework is validated — 18 documented gaps remain before live deployment")

    status_items = [
        (GREEN,  "COMPLETE",  "53/53 agents validated on FSC-2417"),
        (GREEN,  "COMPLETE",  "1,065 automated unit tests — all passing"),
        (GREEN,  "COMPLETE",  "12 gates implemented and verified"),
        (GREEN,  "COMPLETE",  "FCA compliance checks validated for HIGH-FCA story"),
        (GREEN,  "COMPLETE",  "HTML validation report and audit output generated"),
        (AMBER,  "PENDING",   "PostgreSQL production database provisioned and migrated"),
        (AMBER,  "PENDING",   "Anthropic API key added to production secrets manager"),
        (AMBER,  "PENDING",   "Jira production OAuth credentials configured"),
        (AMBER,  "PENDING",   "Copado webhook registered for CI/CD pipeline integration"),
        (AMBER,  "PENDING",   "FCA Evidence Pack reviewed by Compliance before first live sprint"),
    ]

    y = 1.6
    for colour, status, text in status_items:
        rect(sl, 0.4, y, 12.5, 0.44, WHITE)
        rect(sl, 0.4, y, 0.22, 0.44, colour)
        textbox(sl, status, 0.72, y + 0.07, 1.3, 0.32,
                size=11, bold=True, color=colour, align=PP_ALIGN.LEFT)
        textbox(sl, text,   2.15, y + 0.07, 10.6, 0.32,
                size=13, color=DARK, align=PP_ALIGN.LEFT)
        y += 0.48

    textbox(sl, "Full gap register with effort estimates and owners in PRODUCTION_READINESS.md",
            0.4, 6.85, 12.5, 0.4, size=12, color=DARK)


# ── SLIDE 10 — Next Steps ─────────────────────────────────────────────────────

def slide_next_steps(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, NAVY)

    textbox(sl, "Recommended Next Steps",
            0.5, 0.4, 12.0, 0.75, size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    rect(sl, 0.5, 1.2, 12.3, 0.05, TEAL)

    actions = [
        ("01", "Approve Production Deployment",
         "Commission infrastructure provisioning (PostgreSQL, secrets, Jira OAuth, Copado webhook).\n"
         "Assign a delivery lead to own the 18-gap closure plan.\n"
         "Target: framework live in production for next quarter sprint cycle."),
        ("02", "Run a Pilot Sprint",
         "Select 3–5 stories from the next sprint backlog.\n"
         "Run the framework alongside the existing manual QE process.\n"
         "Measure: AC quality delta, FCA scenario coverage, time saved."),
        ("03", "Brief Compliance & FCA Readiness Team",
         "Share the FCA Evidence Pack output from FSC-2417 with the Compliance team.\n"
         "Confirm audit trail format meets SYSC and Consumer Duty documentation requirements.\n"
         "Obtain sign-off that automated evidence packs are acceptable for FCA review."),
    ]

    y = 1.5
    for num, title, body in actions:
        rect(sl, 0.5, y, 12.3, 1.7, DARK)
        rect(sl, 0.5, y, 0.65, 1.7, TEAL)
        textbox(sl, num,   0.55, y + 0.55, 0.6, 0.6, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        textbox(sl, title, 1.35, y + 0.1,  11.3, 0.55, size=18, bold=True, color=WHITE)
        multiline(sl, body.split("\n"), 1.35, y + 0.65, 11.1, 0.95, size=13,
                  color=RGBColor(0xCC, 0xDD, 0xEE))
        y += 1.85

    textbox(sl, "Prepared by QE Architecture Team  ·  FSC Agentic QE Framework  ·  May 2026  ·  CONFIDENTIAL",
            0.5, 7.1, 12.3, 0.35, size=11, color=RGBColor(0x66, 0x77, 0x88), align=PP_ALIGN.LEFT)


# ── Build ─────────────────────────────────────────────────────────────────────

def build():
    prs = _prs()
    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_pipeline(prs)
    slide_capabilities(prs)
    slide_fca(prs)
    slide_metrics(prs)
    slide_benefits(prs)
    slide_production(prs)
    slide_next_steps(prs)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}  ({OUT_PATH.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    build()
